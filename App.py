import streamlit as st
import google.generativeai as genai
import PyPDF2
import io
import json
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- 1. CONFIGURATION ---
st.set_page_config(page_title="Strategic AI Analyst (Pro)", layout="wide", page_icon="📊")
st.title("📊 Strategic AI Analyst: The 'Density' Engine")
st.markdown("Generates a **high-density** strategy deck with KPI Dashboards, Tables, and detailed commentary.")

# --- 2. AUTOMATIC KEY (No Sidebar Needed) ---
api_key = st.secrets["GOOGLE_API_KEY"]

# --- 3. HELPER FUNCTIONS ---

def get_pdf_text(uploaded_file):
    reader = PyPDF2.PdfReader(io.BytesIO(uploaded_file.read()))
    text = ""
    # Read first 50 pages
    for i in range(min(50, len(reader.pages))):
        text += reader.pages[i].extract_text()
    return text

def create_dense_deck(data, chart_buffer):
    prs = Presentation()
    
    # --- SLIDE 1: TITLE SLIDE ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"Strategic Due Diligence: {data.get('company_name', 'Unknown Company')}"
    slide.placeholders[1].text = f"FY{data.get('current_fy', '2024')} Deep Dive Analysis\nStrictly Confidential"

    # --- SLIDE 2: EXECUTIVE SUMMARY (With KPI Sidebar) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank Layout
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    title.text_frame.text = "Executive Summary: Strategic Posture & Key Metrics"
    title.text_frame.paragraphs[0].font.size = Pt(24)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102) # Navy Blue

    # Main Content (Left Side)
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.5), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True # Ensure text wraps
    
    # Intro
    p = tf.add_paragraph()
    p.text = data.get('exec_summary_intro', 'Summary not available.')
    p.font.bold = True
    p.font.size = Pt(14)
    p.space_after = Pt(14)
    
    # Highlights
    for point in data.get('highlights', []):
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(12)
        p.space_after = Pt(10)
        
    # KPI Sidebar (Right Side - Dark Blue Box)
    sidebar = slide.shapes.add_shape(
        1, # Rectangle
        Inches(7.2), Inches(1.5), Inches(2.3), Inches(5)
    )
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = RGBColor(0, 51, 102) # Navy
    sidebar.line.color.rgb = RGBColor(0, 51, 102)
    
    # Add KPIs inside Sidebar
    tf_side = sidebar.text_frame
    tf_side.word_wrap = True
    p = tf_side.paragraphs[0]
    p.text = "KEY METRICS"
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Safe KPI Extraction
    kpis = data.get('kpis', {})
    metrics = [
        ("EBITDA Margin", kpis.get('ebitda_margin', 'N/A')),
        ("ROE", kpis.get('roe', 'N/A')),
        ("Net Debt/EBITDA", kpis.get('leverage', 'N/A'))
    ]
    
    for label, value in metrics:
        p = tf_side.add_paragraph()
        p.text = "\n" + str(value)
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 0) # Yellow
        p.alignment = PP_ALIGN.CENTER
        
        p = tf_side.add_paragraph()
        p.text = label
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER

    # --- SLIDE 3: FINANCIAL DEEP DIVE (Chart + Data Table) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank Layout
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    title.text_frame.text = "Financial Performance Trajectory"
    title.text_frame.paragraphs[0].font.size = Pt(24)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # 1. The Chart (Left Side)
    slide.shapes.add_picture(chart_buffer, Inches(0.5), Inches(1.2), Inches(5.5), Inches(3.5))
    
    # 2. The Data Table (Below Chart)
    financials = data.get('financials', [])
    r, c = len(financials) + 1, 3 
    table_shape = slide.shapes.add_table(r, c, Inches(0.5), Inches(5.0), Inches(5.5), Inches(1.5))
    table = table_shape.table
    
    # Headers
    headers = ["Year", "Revenue (Cr)", "PAT (Cr)"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.text_frame.paragraphs[0].font.size = Pt(11)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(220, 220, 220) # Light Grey
        
    # Fill Data
    for i, fin in enumerate(financials):
        row = i + 1
        table.cell(row, 0).text = str(fin.get('year', '-'))
        table.cell(row, 1).text = str(fin.get('revenue', '-'))
        table.cell(row, 2).text = str(fin.get('net_profit', '-'))
        for j in range(3):
            table.cell(row, j).text_frame.paragraphs[0].font.size = Pt(11)

    # 3. Commentary (Right Side)
    textbox = slide.shapes.add_textbox(Inches(6.1), Inches(1.2), Inches(3.4), Inches(5))
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    
    p = tf.paragraphs[0]
    p.text = "Analyst Commentary:"
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.font.size = Pt(14)
    p.space_after = Pt(12)
    
    for insight in data.get('financial_insights', []):
        p = tf.add_paragraph()
        p.text = f"• {insight}" 
        p.font.size = Pt(11)
        p.space_after = Pt(10)

    # --- SLIDE 4: RISK & MITIGATION MATRIX ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Risk Matrix & Mitigation Strategies"
    
    tf = slide.placeholders[1].text_frame
    tf.clear()
    
    # CRASH PROOF LOOP
    for risk in data.get('risks', []):
        # Risk Title (Safe Get)
        risk_type = risk.get('type', 'General Risk')
        risk_desc = risk.get('risk', 'Risk description unavailable')
        
        p = tf.add_paragraph()
        p.text = f"⚠️ {risk_type}: {risk_desc}"
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(192, 0, 0) # Red
        
        # Mitigation (Safe Get - THIS FIXED YOUR CRASH)
        mitigation = risk.get('mitigation', 'Mitigation strategy not explicitly stated in report.')
        
        p = tf.add_paragraph()
        p.text = f"   Mitigation/Impact: {mitigation}"
        p.font.italic = True
        p.font.size = Pt(12)
        p.space_after = Pt(10)

    # Save
    binary_output = io.BytesIO()
    prs.save(binary_output)
    binary_output.seek(0)
    return binary_output

def generate_chart(data_json):
    years = [item.get('year', 'N/A') for item in data_json]
    revenue = [item.get('revenue', 0) for item in data_json]
    profit = [item.get('net_profit', 0) for item in data_json]
    
    # Reverse to show Oldest -> Newest
    years = years[::-1] 
    revenue = revenue[::-1]
    profit = profit[::-1]

    plt.style.use('bmh') 
    fig, ax = plt.subplots(figsize=(8, 5))
    x = range(len(years))
    
    ax.bar([i - 0.2 for i in x], revenue, width=0.4, label='Revenue', color='#003366')
    ax.bar([i + 0.2 for i in x], profit, width=0.4, label='Net Profit', color='#DAA520')
    
    ax.set_xticks(x)
    ax.set_xticklabels(years)
    ax.legend()
    
    img_buffer = io.BytesIO()
    plt.savefig(img_buffer, format='png', bbox_inches='tight')
    img_buffer.seek(0)
    return img_buffer

# --- 4. MAIN LOGIC ---
uploaded_file = st.file_uploader("Upload Annual Report (PDF)", type="pdf")

if uploaded_file and api_key:
    if st.button("🚀 Generate Dense Deck"):
        
        with st.spinner("Step 1: Ingesting 50 Pages..."):
            text = get_pdf_text(uploaded_file)
            
        with st.spinner("Step 2: Extracting KPIs & Deep Insights..."):
            try:
                genai.configure(api_key=api_key)
                model = genai.GenerativeModel('gemini-flash-latest') 
                
                # --- THE DENSITY PROMPT ---
                prompt = f"""
                Act as a Senior Partner at a top consulting firm.
                Analyze the provided Annual Report text.
                
                Goal: Extract data for a 'High Density' strategy deck.
                
                CRITICAL: 
                1. Identify the Company Name dynamically.
                2. If a specific metric (like ROE) is not found, estimate it from context or put "N/A".
                3. Risks must include a 'Mitigation' or 'Impact' statement.
                
                Output strict JSON:
                {{
                    "company_name": "EXTRACTED_NAME",
                    "current_fy": "2024",
                    "exec_summary_intro": "A 40-word robust strategic summary of the fiscal year.",
                    "kpis": {{
                        "ebitda_margin": "15.5%",
                        "roe": "12%",
                        "leverage": "0.5x"
                    }},
                    "highlights": [
                        "Detailed point 1 (20 words) with specific numbers/growth %.",
                        "Detailed point 2 (20 words) regarding new product launches.",
                        "Detailed point 3 (20 words) regarding geography/market expansion.",
                        "Detailed point 4 (20 words) regarding capex or investment."
                    ],
                    "financials": [
                        {{"year": "2024", "revenue": 1000, "net_profit": 100}},
                        {{"year": "2023", "revenue": 900, "net_profit": 90}},
                        {{"year": "2022", "revenue": 800, "net_profit": 80}}
                    ],
                    "financial_insights": [
                        "Insight on Revenue Growth Drivers (Volume vs Price).",
                        "Insight on Cost Structure or Margin changes.",
                        "Insight on Capital Allocation or Debt reduction."
                    ],
                    "risks": [
                        {{"type": "Operational", "risk": "Supply chain disruption...", "mitigation": "Diversifying suppliers..."}},
                        {{"type": "Market", "risk": "EV slowdown...", "mitigation": "Focusing on Hybrid..."}},
                        {{"type": "Regulatory", "risk": "Carbon tax...", "mitigation": "Green manufacturing..."}}
                    ]
                }}
                
                Source Text: {text}
                """
                
                response = model.generate_content(prompt)
                clean_json = response.text.replace("```json", "").replace("```", "")
                data = json.loads(clean_json)
                
                st.success(f"Analysis Complete: {data.get('company_name', 'Company')}")
                
            except Exception as e:
                st.error(f"AI Error: {e}")
                st.stop()
                
        with st.spinner("Step 3: Rendering Visuals (Tables, Sidebars)..."):
            chart = generate_chart(data.get('financials', []))
            deck = create_dense_deck(data, chart)
            
            clean_name = data.get('company_name', 'Strategy').replace(" ", "_")
            st.download_button(
                label=f"📥 Download Dense Deck: {clean_name}",
                data=deck,
                file_name=f"{clean_name}_Strategy_Deck.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"

            )
